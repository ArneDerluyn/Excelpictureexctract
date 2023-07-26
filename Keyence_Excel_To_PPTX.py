#Importing the modules
import openpyxl
from openpyxl_image_loader import SheetImageLoader
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Cm
import os

# Class to hold 4 images and the title, and create slide
class Microscope:
    def __init__(self,cells_tuple, par_dir):
        self.cells = cells_tuple  # position in excel for the data

        self.name = str(cells_tuple[0].value)  # sample name
        self.parent_directory = par_dir  # Top folder for project


        self.image1 = None
        self.image2 = None
        self.image3 = None
        self.image4 = None

        self.imagename_1 = str(self.name + '_1' + '.png')
        self.imagename_2 = str(self.name + '_2' + '.png')
        self.imagename_3 = str(self.name + '_3' + '.png')
        self.imagename_4 = str(self.name + '_4' + '.png')
    def load_images(self):
        image_loader = SheetImageLoader(sheet)
        self.image1 = image_loader.get(self.cells[1].coordinate)
        self.image2 = image_loader.get(self.cells[2].coordinate)
        self.image3 = image_loader.get(self.cells[3].coordinate)
        self.image4 = image_loader.get(self.cells[4].coordinate)

    def save_images(self):
        cur_path = os.path.join(self.parent_directory, self.name)  # join parent directory and current name
        try:
            os.mkdir(cur_path)  # create folder for microscope images
        except OSError:
            print('Images of ' + self.name + ' will not be saved, slide will not be made')

        os.chdir(cur_path)

        self.image1.save(self.imagename_1)
        self.image2.save(self.imagename_2)
        self.image3.save(self.imagename_3)
        self.image4.save(self.imagename_4)

        os.chdir(self.parent_directory)

    def create_slide(self, prs):
        # Go to current directory
        cur_path = os.path.join(self.parent_directory, self.name)
        os.chdir(cur_path)

        # select layout for slide
        lay1 = prs.slide_layouts[5]

        #create slide
        slide = prs.slides.add_slide(lay1)

        shapes = slide.shapes
        shapes.title.text = self.name

        left = Cm(1)
        top = Cm(5)
        width = Cm(5)
        inbetween = Cm(1)

        pic1 = slide.shapes.add_picture(self.imagename_1, left, top, width)

        left = left + width + inbetween
        pic2 = slide.shapes.add_picture(self.imagename_2, left, top, width)

        left = left + width + inbetween
        pic3 = slide.shapes.add_picture(self.imagename_3, left, top, width)

        left = left + width + inbetween
        pic4 = slide.shapes.add_picture(self.imagename_4, left, top, width)

        os.chdir(self.parent_directory)
        return prs





file_path = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx")])  # Select excel file
par_directory = file_path.rsplit('/', 1)[0]  # Slice selected excel file path to give a working directory to the microscope obkect
filename = file_path.rsplit('/', 1)[1]
filename = file_path.rsplit('.', 1)[0]

pxl_doc = openpyxl.load_workbook(file_path)  # Open Excel File
sheet = pxl_doc.worksheets[0]  # Load first Excel sheet

content = tuple(sheet.rows)  # Tuple containing all the filled cells => iterable to the end

presentation = Presentation()


i = 2
while i < len(content):
    mic = Microscope(content[i],par_directory)
    mic.load_images()
    mic.save_images()
    presentation = mic.create_slide(presentation)
    i = i+1


presentation.save(filename + '.pptx')
